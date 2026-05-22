#!/usr/bin/env python3
"""
Exporta el newsletter como presentación PowerPoint (.pptx)
Genera una diapositiva por sección con las noticias del día.
"""
import json
import os
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Paleta de colores
DARK_BLUE   = RGBColor(0x1A, 0x23, 0x7E)
MID_BLUE    = RGBColor(0x28, 0x35, 0x93)
ACCENT_TEAL = RGBColor(0x00, 0x89, 0x7B)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xF5, 0xF5, 0xF5)
TEXT_DARK   = RGBColor(0x21, 0x21, 0x21)
TEXT_GRAY   = RGBColor(0x75, 0x75, 0x75)
ALERT_RED   = RGBColor(0xC6, 0x28, 0x28)
ALERT_AMBER = RGBColor(0xE6, 0x5C, 0x00)
ALERT_GREEN = RGBColor(0x2E, 0x7D, 0x32)

SECTIONS = [
    ('geopolitica',              '🌍', 'Geopolítica',              DARK_BLUE),
    ('economia_global',          '💹', 'Economía Global',          MID_BLUE),
    ('economia_chile',           '🇨🇱', 'Chile Estratégico',       RGBColor(0x0D, 0x47, 0xA1)),
    ('cooperativismo',           '🤝', 'Cooperativismo',           RGBColor(0x00, 0x60, 0x64)),
    ('cmf',                      '🏛', 'CMF / Regulación',         RGBColor(0x31, 0x18, 0x7E)),
    ('noticias_economicas_actuales', '📈', 'Mercados',             RGBColor(0x1B, 0x5E, 0x20)),
    ('tendencias_tech',          '🚀', 'Tendencias Tech',          RGBColor(0x4A, 0x14, 0x8C)),
    ('inteligencia_artificial',  '🤖', 'Inteligencia Artificial',  RGBColor(0x00, 0x4D, 0x40)),
]

IMPACT_KEYWORDS = [
    'crisis', 'colapso', 'emergencia', 'histórico', 'sanciones',
    'guerra', 'china', 'fed', 'bce', 'reforma tributaria', 'cobre',
    'litio', 'regulación', 'record', 'máximo', 'mínimo',
]


def score_news(title, summary=''):
    text = (title + ' ' + summary).lower()
    score = 5
    for kw in IMPACT_KEYWORDS:
        if kw in text:
            score += 2
    return min(score, 10)


def impact_color(score):
    if score >= 9:
        return ALERT_RED
    if score >= 7:
        return ALERT_AMBER
    return ALERT_GREEN


def add_slide_background(slide, color):
    """Fill slide background with solid color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text,
                font_size=12, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


# ── Slide 1: Portada ────────────────────────────────────────────────────────

def build_cover(prs, date_str):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    add_slide_background(slide, DARK_BLUE)

    W = prs.slide_width
    H = prs.slide_height

    # Franja decorativa superior
    add_rect(slide, 0, 0, W, Inches(0.12), ACCENT_TEAL)

    # Título principal
    add_textbox(slide, Inches(1), Inches(1.8), Inches(8), Inches(1.2),
                'Newsletter Estratégico Financiero',
                font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Subtítulo
    add_textbox(slide, Inches(1), Inches(3.1), Inches(8), Inches(0.6),
                'Análisis diario · Geopolítica · Economía · Mercados · IA',
                font_size=16, color=RGBColor(0xBB, 0xC7, 0xFF), align=PP_ALIGN.CENTER)

    # Fecha
    add_textbox(slide, Inches(1), Inches(3.9), Inches(8), Inches(0.5),
                date_str, font_size=14, bold=True,
                color=ACCENT_TEAL, align=PP_ALIGN.CENTER)

    # Franja inferior
    add_rect(slide, 0, H - Inches(0.9), W, Inches(0.9), MID_BLUE)
    add_textbox(slide, Inches(0.5), H - Inches(0.8), Inches(9), Inches(0.6),
                'erickfranciscohernandez.github.io/financiero',
                font_size=11, color=RGBColor(0xBB, 0xC7, 0xFF), align=PP_ALIGN.CENTER)


# ── Slide 2: Resumen ejecutivo ───────────────────────────────────────────────

def build_summary(prs, all_news):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_slide_background(slide, LIGHT_GRAY)

    W = prs.slide_width

    # Header
    add_rect(slide, 0, 0, W, Inches(1.0), DARK_BLUE)
    add_textbox(slide, Inches(0.4), Inches(0.2), Inches(9), Inches(0.6),
                '📋  Resumen Ejecutivo', font_size=22, bold=True, color=WHITE)

    # Top noticias con mayor score
    top_news = []
    for section_key, _, label, _ in SECTIONS:
        for item in all_news.get(section_key, []):
            s = score_news(item.get('title', ''), item.get('summary', ''))
            top_news.append((s, label, item.get('title', ''), item.get('source', '')))

    top_news.sort(reverse=True)
    top_news = top_news[:8]

    y = Inches(1.15)
    row_h = Inches(0.52)
    for i, (sc, label, title, source) in enumerate(top_news):
        bg = WHITE if i % 2 == 0 else RGBColor(0xEE, 0xF1, 0xFB)
        add_rect(slide, Inches(0.3), y, W - Inches(0.6), row_h, bg)

        # Score badge
        badge_color = impact_color(sc)
        add_rect(slide, Inches(0.3), y, Inches(0.55), row_h, badge_color)
        add_textbox(slide, Inches(0.3), y + Pt(6), Inches(0.55), row_h - Pt(6),
                    str(sc), font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # Sección
        add_textbox(slide, Inches(0.95), y + Pt(4), Inches(1.4), Pt(14),
                    label, font_size=9, bold=True, color=DARK_BLUE)

        # Título
        add_textbox(slide, Inches(0.95), y + Pt(18), Inches(7.5), Pt(20),
                    title[:95] + ('…' if len(title) > 95 else ''),
                    font_size=10, color=TEXT_DARK)

        y += row_h + Inches(0.04)


# ── Slide por sección ────────────────────────────────────────────────────────

def build_section_slide(prs, section_key, icon, label, header_color, news_list):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_slide_background(slide, LIGHT_GRAY)

    W = prs.slide_width

    # Header
    add_rect(slide, 0, 0, W, Inches(1.0), header_color)
    add_textbox(slide, Inches(0.4), Inches(0.22), Inches(9), Inches(0.65),
                f'{icon}  {label}', font_size=24, bold=True, color=WHITE)

    if not news_list:
        add_textbox(slide, Inches(0.5), Inches(1.5), Inches(9), Inches(1),
                    'Sin noticias disponibles para esta categoría.',
                    font_size=13, color=TEXT_GRAY)
        return

    y = Inches(1.1)
    card_gap = Inches(0.12)

    for item in news_list[:4]:
        title   = item.get('title', '')
        summary = item.get('summary', '')[:160]
        source  = item.get('source', '')
        sc      = score_news(title, summary)

        card_h = Inches(1.35)
        # Card background
        add_rect(slide, Inches(0.3), y, W - Inches(0.6), card_h, WHITE)
        # Left accent bar
        add_rect(slide, Inches(0.3), y, Inches(0.06), card_h, impact_color(sc))

        # Score badge
        add_rect(slide, W - Inches(1.1), y + Inches(0.1),
                 Inches(0.7), Inches(0.3), impact_color(sc))
        add_textbox(slide, W - Inches(1.1), y + Inches(0.1),
                    Inches(0.7), Inches(0.3),
                    f'{sc}/10', font_size=9, bold=True,
                    color=WHITE, align=PP_ALIGN.CENTER)

        # Title
        add_textbox(slide, Inches(0.5), y + Inches(0.1),
                    W - Inches(1.8), Inches(0.38),
                    title[:90] + ('…' if len(title) > 90 else ''),
                    font_size=11, bold=True, color=TEXT_DARK)

        # Summary
        add_textbox(slide, Inches(0.5), y + Inches(0.48),
                    W - Inches(0.9), Inches(0.55),
                    summary + ('…' if len(item.get('summary', '')) > 160 else ''),
                    font_size=9, color=TEXT_GRAY)

        # Source
        add_textbox(slide, Inches(0.5), y + Inches(1.05),
                    Inches(4), Inches(0.25),
                    f'📰 {source}', font_size=8, color=TEXT_GRAY)

        y += card_h + card_gap


# ── Slide final ──────────────────────────────────────────────────────────────

def build_closing(prs, date_str, total_news):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    add_slide_background(slide, DARK_BLUE)

    W = prs.slide_width
    H = prs.slide_height

    add_rect(slide, 0, 0, W, Inches(0.12), ACCENT_TEAL)

    add_textbox(slide, Inches(1), Inches(1.5), Inches(8), Inches(0.8),
                '📊', font_size=48, align=PP_ALIGN.CENTER, color=WHITE)

    add_textbox(slide, Inches(1), Inches(2.6), Inches(8), Inches(0.7),
                f'{total_news} noticias analizadas',
                font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1), Inches(3.4), Inches(8), Inches(0.5),
                f'Edición {date_str}',
                font_size=14, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1), Inches(4.1), Inches(8), Inches(0.5),
                'erickfranciscohernandez.github.io/financiero',
                font_size=12, color=RGBColor(0xBB, 0xC7, 0xFF), align=PP_ALIGN.CENTER)

    add_rect(slide, 0, H - Inches(0.9), W, Inches(0.9), MID_BLUE)
    add_textbox(slide, Inches(0.5), H - Inches(0.75), Inches(9), Inches(0.55),
                'Newsletter Estratégico Financiero · Generado automáticamente',
                font_size=10, color=RGBColor(0xBB, 0xC7, 0xFF), align=PP_ALIGN.CENTER)


# ── Main ─────────────────────────────────────────────────────────────────────

def generate_pptx(output_path=None):
    date_str = datetime.now().strftime('%d de %B de %Y').replace('May', 'Mayo')
    date_file = datetime.now().strftime('%Y%m%d')

    if output_path is None:
        output_path = f'newsletter_{date_file}.pptx'

    # Cargar noticias
    try:
        with open('noticias_diarias.json', 'r', encoding='utf-8') as f:
            all_news = json.load(f)
    except Exception as e:
        print(f'❌ No se pudo leer noticias_diarias.json: {e}')
        return None

    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.63)  # 16:9

    build_cover(prs, date_str)
    build_summary(prs, all_news)

    for section_key, icon, label, color in SECTIONS:
        news_list = all_news.get(section_key, [])
        build_section_slide(prs, section_key, icon, label, color, news_list)

    total = sum(len(all_news.get(k, [])) for k, *_ in SECTIONS)
    build_closing(prs, date_str, total)

    prs.save(output_path)
    print(f'✅ PowerPoint generado: {output_path}')
    print(f'   • {len(prs.slides)} diapositivas')
    print(f'   • {total} noticias incluidas')
    return output_path


if __name__ == '__main__':
    generate_pptx()
